#
# PRO: DSTK.py
#
# AUTHOR: RNH
#
# DATE: 8-6-2021
#
# IMPORT REQUIRED PACKAGES
import os 
import pandas as pd
import numpy as np
import seaborn as sns
import matplotlib
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import shutil

import py_plot

os.chdir(r'D:\Documents\Python\DATA_SCIENCE_TOOLKIT')
print(os.getcwd())

###################
# INPUT VARS HERE #
###################

# IMPORT DATA
df = pd.read_csv('AB_NYC_2019.csv', low_memory= False)
#df = pd.read_excel('OPTIONS_APPL.xlsx', sheet_name = 'jtymg7d6kmxpnhde')
n_row, n_col = df.shape


# FILE NAME
f_name = 'POWERPOINT' + '.pptx'



# How to create a slide 
prs = Presentation()

title_only_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

top = Inches(.25)
left = Inches(.5)
width = Inches(9)
height = Inches(1.25)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

p = tf.add_paragraph()
p.text = "Overview"
p.font.size = Pt(36)
p.alignment = PP_ALIGN.LEFT



def add_slide_continuous_var(df, x_var_name):

    #x_var_name = 'minimum_nights'
    sum_stats = py_plot.box_hist(df, x_var_name)
    #print(sum_stats)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # ADD TITLE TEXT
    top = Inches(.25)
    left = Inches(.5)
    width = Inches(9)
    height = Inches(1.25)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = x_var_name
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.LEFT


    # ADD PHOTO
    left = Inches(0)
    top = Inches(3)
    pic = slide.shapes.add_picture('temp_plots/' + x_var_name + ".png", left, top)


    # ADD TABLE
    #x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)

    left = Inches(6.5)
    top = Inches(2.75)
    width = Inches(3.25)
    height = Inches(4.25)
    # add_table = (rows, cols, left, top, width, height)
    shape = slide.shapes.add_table(7, 2, left, top, width, height)
    table = shape.table

    # Merge the the first row.
    cell = table.cell(0, 0)
    other_cell = table.cell(0, 1)
    cell.merge(other_cell)

    cell = table.cell(0, 0)
    cell.text = '         Summary Statistics'

    # FILL IN THE TABLE
    cell = table.cell(1, 0)
    cell.text = '25th percentile'
    cell.alignment = PP_ALIGN.CENTER
    cell = table.cell(1, 1)
    cell.text = str(sum_stats['25th_p'].iloc[0])

    cell = table.cell(2, 0)
    cell.text = 'Median'
    cell = table.cell(2, 1)
    cell.text = str(sum_stats['median'].iloc[0])


    cell = table.cell(3, 0)
    cell.text = '75th percentile'
    cell = table.cell(3, 1)
    cell.text = str(sum_stats['75th_p'].iloc[0])

    cell = table.cell(4, 0)
    cell.text = 'Mean'
    cell = table.cell(4, 1)
    cell.text = str(sum_stats['mean'].iloc[0])


    cell = table.cell(5, 0)
    cell.text = 'Sample Size'
    cell = table.cell(5, 1)
    cell.text = str(sum_stats['N'].iloc[0])

    cell = table.cell(6, 0)
    cell.text = '% Blank'
    cell = table.cell(6, 1)
    cell.text = str(sum_stats['p_blank'].iloc[0])




def add_slide_cat_var(df, x_var_name):

    #x_var_name = 'minimum_nights'
    sum_stats = py_plot.bar_chart(df, x_var_name)
    #print(sum_stats)
    n_row, n_col = sum_stats.shape

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # ADD TITLE TEXT
    top = Inches(.25)
    left = Inches(.5)
    width = Inches(9)
    height = Inches(1.25)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = x_var_name
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.LEFT


    # ADD PHOTO
    left = Inches(0)
    top = Inches(1.5)
    pic = slide.shapes.add_picture('temp_plots/' + x_var_name + ".png", left, top)


    # ADD TABLE
    #x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)

    left = Inches(5.5)
    top = Inches(.5)
    width = Inches(4.5)
    height = Inches(6.25)
    # add_table = (rows, cols, left, top, width, height)
    shape = slide.shapes.add_table((n_row+1), n_col, left, top, width, height)
    table = shape.table

    c = list(sum_stats.columns)

    for count, value in enumerate(c):
        cell = table.cell(0, count)
        #cell.text = Pt(12)
        cell.text = str(value)
        #cell.text.font.size = Pt(12)
        #p.alignment = PP_ALIGN.LEFT


    # Merge the the first row.
    #cell = table.cell(0, 0)
    #other_cell = table.cell(0, 1)
    #cell.merge(other_cell)
    #cell = table.cell(0, 0)
    #cell.text = '         Summary Statistics'

    # FILL IN THE TABLE
    for i in range(0,n_row):

        for j in range(0, n_col):
            cell = table.cell(i+1, j)
            cell.text = str(sum_stats.iloc[i,j])


#################
# FOR LOOP HERE #
#################

cols = df.columns

for c in cols: 
    t = df.dtypes[c]

    if t == object:
        print(c + "   "+str(t))
        add_slide_cat_var(df, c)
    else:
        u = int(len(df[c].unique()))
        print(c + "   "+str(t))
        try:
            add_slide_continuous_var(df, c)
        except:
            print("An exception occurred")
        

# SAVE THE FINAL RESULTS HERE.
prs.save(f_name)

shutil.rmtree('temp_plots')

os.system(f_name)


print(" --- PROGRAM COMPLETE  -----")



